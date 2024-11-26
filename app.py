from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from dotenv import load_dotenv
import io
import openai
import os

# Cargar variables de entorno desde el archivo .env
load_dotenv()

# Configuración de la clave de API de OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    raise ValueError("La clave de API de OpenAI no está configurada. Por favor, configura la variable de entorno 'OPENAI_API_KEY'.")

# Crea la instancia de Flask
app = Flask(__name__)
CORS(app)  # Configura CORS para permitir solicitudes desde el cliente

# Función para procesar el texto usando GPT de OpenAI
def procesar_con_ia(texto):
    try:
        respuesta = openai.Completion.create(
            engine="text-davinci-003",
            prompt=f"Divide este texto en puntos clave para diapositivas:\n\n{texto}",
            max_tokens=500,
            temperature=0.5
        )
        return respuesta.choices[0].text.strip()
    except Exception as e:
        print(f"Error en la llamada a la API de OpenAI: {e}")
        return None

# Ruta para generar la vista previa
@app.route('/generate-preview', methods=['POST'])
def generate_preview():
    try:
        
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({"error": "Falta el parámetro 'text'"}), 400

        text = data['text'].strip()
        if not text:
            return jsonify({"error": "El texto proporcionado está vacío"}), 400

        processed_text = procesar_con_ia(text)
        if processed_text is None:
            return jsonify({"error": "Error procesando el texto con la IA"}), 500

        slides = processed_text.split("\n\n")
        slide_previews = [{"content": slide.strip()} for slide in slides if slide.strip()]

        if not slide_previews:
            return jsonify({"error": "No se generaron puntos clave para la vista previa"}), 400

        return jsonify({"slides": slide_previews})
    except Exception as e:
        print(f"Error interno: {e}")
        return jsonify({"error": "Error interno del servidor"}), 500

# Ruta para generar y descargar diapositivas
@app.route('/generate-slides', methods=['POST'])
def generate_slides():
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({"error": "Falta el parámetro 'text'"}), 400

        text = data['text'].strip()
        if not text:
            return jsonify({"error": "El texto proporcionado está vacío"}), 400

        processed_text = procesar_con_ia(text)
        if processed_text is None:
            return jsonify({"error": "Error procesando el texto con la IA"}), 500

        slides = processed_text.split("\n\n")

        # Crear la presentación en PowerPoint
        prs = Presentation()
        for slide_content in slides:
            slide_content = slide_content.strip()
            if not slide_content:
                continue
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = "Punto Clave"
                content.text = slide_content
            except Exception as slide_error:
                print(f"Error al crear una diapositiva: {slide_error}")

        # Guarda el archivo PPTX en memoria
        pptx_file = io.BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)

        return send_file(pptx_file, as_attachment=True, download_name="diapositivas.pptx")
    except Exception as e:
        print(f"Error interno: {e}")
        return jsonify({"error": "Error interno del servidor"}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5001)
